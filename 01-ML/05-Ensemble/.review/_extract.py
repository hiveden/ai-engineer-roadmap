#!/usr/bin/env python3
"""提取 Ensemble 章节两个 PPT 的全文（标题、正文、表格、备注）"""
from pptx import Presentation

def extract_pptx(path):
    prs = Presentation(path)
    out = []
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"\n## Slide {i}\n")
        # title
        try:
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                t = slide.shapes.title.text_frame.text.strip()
                if t:
                    out.append(f"**标题**: {t}")
        except Exception:
            pass
        for shape in slide.shapes:
            if shape.has_text_frame:
                # 跳过已捕获的标题
                try:
                    if slide.shapes.title is not None and shape._element is slide.shapes.title._element:
                        continue
                except Exception:
                    pass
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        out.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    out.append(" | ".join(cells))
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                out.append(f"\n_备注_: {notes}")
    return "\n".join(out), len(prs.slides)

base = "/Users/xuelin/projects/ai-engineer-roadmap/01-ML/05-Ensemble"
out_path = f"{base}/.review/PPT-extracted.md"

main_text, n1 = extract_pptx(f"{base}/06集成学习.pptx")
supp_text, n2 = extract_pptx(f"{base}/06集成学习补充材料.pptx")

with open(out_path, "w", encoding="utf-8") as f:
    f.write("# 主 PPT (06集成学习.pptx)\n")
    f.write(main_text)
    f.write("\n\n---\n\n")
    f.write("# 补充材料 (06集成学习补充材料.pptx)\n")
    f.write(supp_text)

print(f"主 PPT slides: {n1}")
print(f"补充材料 slides: {n2}")
print(f"输出: {out_path}")
