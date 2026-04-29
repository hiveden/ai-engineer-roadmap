from pptx import Presentation
import sys
import re

def extract_text(pptx_file):
    try:
        prs = Presentation(pptx_file)
        md_content = f"# {pptx_file.split('/')[-1].replace('.pptx', '')}\n\n"
        
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    # Clean up multiple newlines/spaces
                    text = re.sub(r'\n+', '\n', text)
                    slide_text.append(text)
            
            if slide_text:
                md_content += f"## Slide {i + 1}\n\n"
                md_content += "\n".join(slide_text) + "\n\n"
                md_content += "---\n\n"
            
        with open(pptx_file.replace(".pptx", ".md"), "w", encoding="utf-8") as f:
            f.write(md_content)
        print("Success")
    except Exception as e:
        print(f"Error: {e}")

extract_text("01-ML/03-knn/02KNN算法.pptx")
